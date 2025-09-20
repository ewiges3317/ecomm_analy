import os
import pandas as pd

# ensure python-pptx is available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception as e:
    print("Missing python-pptx. Run:  pip install python-pptx")
    raise

BASE = r"data\processed\analysis"
DOCS = r"docs"
OUT_DIR = r"outputs\reports"
IMG_DIR = r"outputs\dashboards"
os.makedirs(OUT_DIR, exist_ok=True)

# Inputs
journey_csv = os.path.join(BASE, "journey_delivery_csat_repeat.csv")
exec_1p_md  = os.path.join(DOCS, "EXEC_1P_Delivery_CSAT_Repeat.md")
impact_txt  = os.path.join(DOCS, "impact_summary.txt")
prio_txt    = os.path.join(DOCS, "prioritize_fixes_summary.txt")

img1 = os.path.join(IMG_DIR, "late_pct_by_state_top.png")
img2 = os.path.join(IMG_DIR, "avg_review_by_delay_bucket.png")
img3 = os.path.join(IMG_DIR, "repeat_90d_on_time_vs_late.png")

# Load journey stats
df = pd.read_csv(journey_csv, low_memory=False)
orders_n = len(df)
late_pct = float(pd.to_numeric(df["is_late"], errors="coerce").mean() * 100)
avg_on   = float(pd.to_numeric(df["review_score"], errors="coerce")[df["is_late"]==0].mean())
avg_lt   = float(pd.to_numeric(df["review_score"], errors="coerce")[df["is_late"]==1].mean())
rep_on   = float(pd.to_numeric(df["repeat_90d"], errors="coerce")[df["is_late"]==0].mean() * 100)
rep_lt   = float(pd.to_numeric(df["repeat_90d"], errors="coerce")[df["is_late"]==1].mean() * 100)
aov = pd.to_numeric(df.get("price_total", pd.Series(dtype=float)), errors="coerce").dropna().mean()
if pd.isna(aov): aov = 137.75

def read_text(path):
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read().strip()
    return ""

impact_block = read_text(impact_txt)
prio_block   = read_text(prio_txt)

# Build deck
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_and_content  = prs.slide_layouts[1]
blank_layout       = prs.slide_layouts[6]

# Slide 1 — Title
s1 = prs.slides.add_slide(title_slide_layout)
s1.shapes.title.text = "Delivery → CSAT → Repeat Growth Engine"
s1.placeholders[1].text = (
    f"Orders analyzed: {orders_n:,}\n"
    f"Late deliveries: {late_pct:.2f}%\n"
    f"Avg review: on-time {avg_on:.2f} vs late {avg_lt:.2f}\n"
    f"Repeat in 90d: on-time {rep_on:.2f}% vs late {rep_lt:.2f}%\n"
    f"Assumed AOV: ${aov:,.2f}"
)

# Slide 2 — KPIs (charts)
s2 = prs.slides.add_slide(blank_layout)
left, top = Inches(0.5), Inches(0.3)
# title textbox
tb = s2.shapes.add_textbox(left, top, Inches(9), Inches(0.6))
tf = tb.text_frame
tf.text = "Key KPIs"
tf.paragraphs[0].font.size = Pt(28)
# images
y = top + Inches(0.8)
if os.path.exists(img1):
    s2.shapes.add_picture(img1, Inches(0.5), y, width=Inches(3.8))
if os.path.exists(img2):
    s2.shapes.add_picture(img2, Inches(4.4), y, width=Inches(3.8))
y2 = y + Inches(3.0)
if os.path.exists(img3):
    s2.shapes.add_picture(img3, Inches(0.5), y2, width=Inches(3.8))

# Slide 3 — Where to fix first ($ impact + priorities)
s3 = prs.slides.add_slide(title_and_content)
s3.shapes.title.text = "Where to Fix First (Impact & Priorities)"
body = s3.placeholders[1].text_frame
body.clear()
p = body.paragraphs[0]
p.text = "Impact summary:"
p.font.bold = True
for line in (impact_block.splitlines()[:8] or ["(impact summary not found)"]):
    para = body.add_paragraph()
    para.text = line
    para.level = 1
body.add_paragraph().text = ""
p2 = body.add_paragraph()
p2.text = "Prioritization:"
p2.font.bold = True
for line in (prio_block.splitlines()[:12] or ["(prioritization not found)"]):
    para = body.add_paragraph()
    para.text = line
    para.level = 1

# Save
out_pptx = os.path.join(OUT_DIR, "Delivery_CSAT_Repeat_Executive.pptx")
prs.save(out_pptx)
print(f"Wrote deck -> {out_pptx}")
